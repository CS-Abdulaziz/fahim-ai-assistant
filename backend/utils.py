import os
import tempfile
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

def extract_text_from_pptx(file_path):

    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def _read_file_bytes(uploaded_file):
    if isinstance(uploaded_file, (bytes, bytearray)):
        return bytes(uploaded_file)

    if hasattr(uploaded_file, "getvalue"):
        return uploaded_file.getvalue()

    if hasattr(uploaded_file, "read"):
        return uploaded_file.read()

    raise ValueError("Unsupported uploaded file type.")

def _get_filename(uploaded_file, filename=None):
    if filename:
        return filename

    detected_name = getattr(uploaded_file, "filename", None) or getattr(uploaded_file, "name", None)
    if not detected_name:
        raise ValueError("Filename is required.")

    return os.path.basename(detected_name)

def process_uploaded_file(uploaded_file, filename=None):
    file_name = _get_filename(uploaded_file, filename)
    file_bytes = _read_file_bytes(uploaded_file)

    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file_name)[1]) as tmp:
        tmp.write(file_bytes)
        file_path = tmp.name

    try:
        lower_name = file_name.lower()

        if lower_name.endswith(".pdf"):
            loader = PyPDFLoader(file_path)
            documents = loader.load()
        elif lower_name.endswith((".pptx", ".ppt")):
            text_content = extract_text_from_pptx(file_path)
            documents = [Document(page_content=text_content, metadata={"source": file_name})]
        elif lower_name.endswith(".docx"):
            loader = Docx2txtLoader(file_path)
            documents = loader.load()
        else:
            loader = TextLoader(file_path)
            documents = loader.load()

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        return text_splitter.split_documents(documents)
    
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)
