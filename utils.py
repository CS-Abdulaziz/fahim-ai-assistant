import os
import tempfile
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

def extract_text_from_pptx(file_path):

    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def process_uploaded_file(uploaded_file):

    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp:
        
        tmp.write(uploaded_file.getvalue())
        file_path = tmp.name

    try:
        
        if uploaded_file.name.endswith(".pdf"):
            loader = PyPDFLoader(file_path)
            documents = loader.load()
        elif uploaded_file.name.endswith((".pptx", ".ppt")):
            text_content = extract_text_from_pptx(file_path)
            documents = [Document(page_content=text_content, metadata={"source": uploaded_file.name})]
        elif uploaded_file.name.endswith(".docx"):
            loader = Docx2txtLoader(file_path)
            documents = loader.load()
        else:
            loader = TextLoader(file_path)
            documents = loader.load()

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        return text_splitter.split_documents(documents)
    
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)